"""Professional PPTX taqdimot yasash uchun tayyor maketlar.

Bu modul sandbox ichiga avtomatik ko'chiriladi, shuning uchun GPT yozgan
kodda to'g'ridan-to'g'ri `import deck` qilish mumkin.

NIMA UCHUN KERAK (docgen bilan bir xil sabab): model shakllarni qo'lda
Inches(...) bilan joylashtirganda ularning HAQIQIY o'lchamini hisoblamaydi.
Natijada rasm matn ustiga chiqadi, manba yozuvi rasm ustiga tushadi,
altbet raqami rasm bilan ustma-ust bo'ladi va slayd "dabdala" ko'rinadi.
Aynan shu uchta xato eng ko'p uchraydi.

Bu yerda GEOMETRIYA MODULGA tegishli: har bir maket o'z ichida xavfsiz
chekka (margin), altbet uchun band joy va rasm uchun aniq quti hisoblaydi.
Modelga faqat MAZMUN qoladi — nima yozish, qaysi rasmni qo'yish.

ODDIY ISHLATISH:

    import deck
    d = deck.Deck("Birinchi jahon urushi", theme="navy")
    d.cover("BIRINCHI JAHON URUSHI", "1914–1918: sabablar va oqibatlar",
            image="rasm1.jpg", credit="iwm.org.uk", footer="Tarix fanidan")
    d.section("Urush sabablari", 1)
    d.bullets("Asosiy sabablar",
              [("Imperializm", "Mustamlakalar uchun kurash"),
               ("Militarizm", "Qurollanish poygasi")],
              image="rasm2.jpg", credit="wikimedia.org")
    d.stats("Raqamlarda", [("38 mln", "Umumiy talafot"), ("4", "Urush yillari")])
    d.image_slide("Xandaq urushi", "rasm3.jpg", "G'arbiy front, 1916",
                  credit="iwm.org.uk")
    d.closing("Xulosa", "Urush jahon tartibini butunlay o'zgartirdi")
    d.save("output/taqdimot.pptx")
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

# 16:9 — zamonaviy proyektor va noutbuk ekrani standarti.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Xavfsiz chekka: hech qanday matn bu chegaradan tashqariga chiqmaydi.
MARGIN = Inches(0.75)
# Altbet uchun BAND joy — hech bir maket bu bandga kontent qo'ymaydi.
FOOTER_H = Inches(0.55)

# Shrift: Calibri Office bilan birga Windows'da ham, Mac'da ham keladi.
# "Segoe UI" chiroyliroq, lekin Mac'da yo'q va almashtirilib ketadi.
FONT = "Calibri"

THEMES = {
    "navy":   {"bg": "0B1E36", "bg2": "163A5F", "fg": "FFFFFF",
               "muted": "A8BDD4", "accent": "F5A623", "card": "12263F"},
    "forest": {"bg": "0E2A22", "bg2": "1B4D3E", "fg": "FFFFFF",
               "muted": "A7C4B5", "accent": "E8B54A", "card": "143529"},
    "plum":   {"bg": "241432", "bg2": "45255C", "fg": "FFFFFF",
               "muted": "C4AFD6", "accent": "F2708B", "card": "2E1A40"},
    "slate":  {"bg": "1C2128", "bg2": "343B45", "fg": "FFFFFF",
               "muted": "AFB8C1", "accent": "58A6FF", "card": "22282F"},
}


# `image=AUTO` — "o'zing tanla": Deck'ga berilgan rasmlar hovuzidan
# navbatdagisi olinadi. `image=None` esa "bu slaydda rasm KERAK EMAS"
# degani. Ikkalasini ajratish shart, aks holda modelga har bir slaydda
# rasm nomini qo'lda yozdirish kerak bo'lardi — va u ko'pincha faqat
# birinchi slaydga yozib, qolganini bo'sh qoldirardi.
AUTO = object()


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _alpha(shape, percent):
    """Shakl to'ldirishiga shaffoflik qo'shadi (0-100).

    python-pptx'da shaffoflik uchun API yo'q — XML'ga qo'lda qo'shiladi.
    Fon rasm ustidagi "scrim" (qorayituvchi qatlam) aynan shu bilan
    ishlaydi: rasm ko'rinib turadi, matn esa o'qiladigan bo'lib qoladi.
    """
    try:
        solid = shape.fill._xPr.find(qn("a:solidFill"))
        srgb = solid.find(qn("a:srgbClr"))
        srgb.append(parse_xml('<a:alpha %s val="%d"/>'
                              % (nsdecls("a"), int(percent) * 1000)))
    except Exception:
        pass          # shaffoflik — bezak, uning yo'qligi slaydni buzmaydi


def _no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def _para(tf, text, size, color, bold=False, space_after=6,
          align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = str(text)
    p.alignment = align
    p.space_after = Pt(space_after)
    f = p.font
    f.name, f.size, f.bold = FONT, Pt(size), bold
    f.color.rgb = _rgb(color)
    return p


def _fit_box(img_path, box_l, box_t, box_w, box_h):
    """Rasmni qutiga NISBATNI BUZMASDAN joylaydi (contain + markazlash).

    Model ko'pincha width va height'ni BIRGA beradi — shunda rasm cho'zilib
    havaskorona ko'rinadi. Bu funksiya har doim bitta o'lchovni tanlaydi.
    """
    try:
        from PIL import Image
        with Image.open(img_path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 4, 3
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return (int(box_l + (box_w - w) / 2), int(box_t + (box_h - h) / 2), w, h)


def _round_corners(pic, radius=12000):
    """Rasm burchaklarini yumaloqlaydi — matn kartochkalari bilan bir xil.

    python-pptx'da bu uchun API yo'q: rasmning geometriyasi XML'da
    to'g'ridan-to'g'ri almashtiriladi. Xato bo'lsa jimgina o'tkazib
    yuboriladi — burchak bezak, rasm esa undan muhimroq.
    """
    try:
        spPr = pic._element.spPr
        old = spPr.find(qn("a:prstGeom"))
        if old is not None:
            spPr.remove(old)
        spPr.append(parse_xml(
            '<a:prstGeom %s prst="roundRect"><a:avLst>'
            '<a:gd name="adj" fmla="val %d"/></a:avLst></a:prstGeom>'
            % (nsdecls("a"), radius)))
    except Exception:
        pass


def _cover_crop(img_path, target_w, target_h):
    """Rasmni quti nisbatiga qirqadi (fon uchun — bo'sh joy qolmasin).

    Qaytaradi: yangi fayl yo'li (yoki xato bo'lsa asl yo'l).
    """
    try:
        from PIL import Image
        with Image.open(img_path) as im:
            im = im.convert("RGB")
            iw, ih = im.size
            target = target_w / target_h
            if iw / ih > target:                 # juda keng — yon tomonlarni qirqamiz
                new_w = int(ih * target)
                left = (iw - new_w) // 2
                im = im.crop((left, 0, left + new_w, ih))
            else:                                # juda baland — usti/ostini qirqamiz
                new_h = int(iw / target)
                top = (ih - new_h) // 3          # 1/3 — odam boshi kesilmasin
                im = im.crop((0, top, iw, top + new_h))
            out = os.path.splitext(img_path)[0] + "_bg.jpg"
            im.save(out, "JPEG", quality=88)
            return out
    except Exception:
        return img_path


class Deck:
    """Taqdimot quruvchi. Har bir metod bitta tugallangan slayd qo'shadi."""

    def __init__(self, title="Taqdimot", theme="navy", footer=None, images=None):
        """`images` — mavjud rasm fayllari ro'yxati (masalan yuklab olingan
        rasm1.jpg, rasm2.jpg ...). Maketlar ulardan O'ZI oladi, ya'ni rasm
        butun taqdimot bo'ylab tarqaladi, faqat muqovada qolib ketmaydi.
        """
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SLIDE_W, SLIDE_H
        self.t = THEMES.get(theme, THEMES["navy"])
        self.title = title
        self.footer_text = footer or title
        self.n = 0
        self._pool = [p for p in (images or []) if p and os.path.exists(p)]
        self._taken = 0

    # ── ichki yordamchilar ──────────────────────────────────────────
    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _pick(self, image):
        """image=AUTO -> hovuzdan navbatdagisi; None -> rasmsiz; yo'l -> o'zi."""
        if image is None:
            return None
        if image is not AUTO:
            return image if os.path.exists(image) else None
        if not self._pool:
            return None
        # Avval ishlatilmaganlari; tugasa boshidan aylanadi — takroriy rasm
        # bo'sh slayddan yaxshiroq, lekin hovuz yetarli bo'lsa takror yo'q.
        path = self._pool[self._taken % len(self._pool)]
        self._taken += 1
        return path

    def _photo(self, slide, image, left, top, width, height, mode="fill"):
        """Rasmni ANIQ ramkaga joylaydi — matn kartochkasi bilan bir xil.

        mode="fill" — rasm ramkani to'liq to'ldiradi (qirqiladi). Yon
        ustundagi rasm uchun: matn kartochkasi bilan bir xil balandlik va
        bir xil chekka beradi, ya'ni ikkala ustun ko'zga bir xil ko'rinadi.

        mode="fit" — rasm butunligicha ko'rinadi, ramka ichida markazda.
        Asosiy fotosurat va diagramma uchun (qirqish mazmunni yo'qotadi).

        ⚠️ "fit" rejimida ramka RASMGA YOPISHADI: berilgan quti faqat
        CHEGARA, ramka esa rasmning haqiqiy o'lchamiga qarab kichrayadi.
        Busiz kvadrat yoki tik rasm keng ramka ichida kichkina bo'lib
        qolib, yon tomonlarda katta bo'sh joy qolardi.

        Qaytaradi: ramkaning haqiqiy (left, top, width, height) —
        chaqiruvchi izoh va manbani aynan uning ostiga qo'yishi uchun.
        """
        exists = bool(image) and os.path.exists(image)
        pad = Inches(0.14)
        # `!= "fill"` (== "fit" emas): noma'lum rejim berilsa ham l/t/w/h
        # aniqlanmay qolmaydi — pastda ular ishlatiladi.
        if mode != "fill" and exists:
            l, t, w, h = _fit_box(image, left + pad, top + pad,
                                  width - 2 * pad, height - 2 * pad)
            rect = (l - pad, t - pad, w + 2 * pad, h + 2 * pad)
        else:
            rect = (left, top, width, height)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, *rect)
        _no_line(card)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(self.t["card"])
        if not exists:
            return rect

        if mode == "fill":
            path = _cover_crop(image, width, height)
            pic = slide.shapes.add_picture(path, left, top, width, height)
        else:
            pic = slide.shapes.add_picture(image, l, t, w, h)
        _round_corners(pic)
        return rect

    def _bg_solid(self, slide, color=None, gradient=True):
        """Tekis yoki gradientli fon. Gradient tekis rangdan ancha jonli."""
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _no_line(r)
        if gradient:
            try:
                r.fill.gradient()
                r.fill.gradient_angle = 45.0
                r.fill.gradient_stops[0].color.rgb = _rgb(color or self.t["bg"])
                r.fill.gradient_stops[1].color.rgb = _rgb(self.t["bg2"])
            except Exception:
                r.fill.solid()
                r.fill.fore_color.rgb = _rgb(color or self.t["bg"])
        else:
            r.fill.solid()
            r.fill.fore_color.rgb = _rgb(color or self.t["bg"])
        return r

    def _bg_photo(self, slide, image, scrim=72):
        """To'liq ekranli fon rasmi + qorayituvchi qatlam.

        Scrim SHART: fotosurat ustidagi oq matn usiz o'qilmaydi. 72% —
        rasm ko'rinib turadigan, lekin matn to'liq o'qiladigan daraja.
        """
        path = _cover_crop(image, SLIDE_W, SLIDE_H)
        slide.shapes.add_picture(path, 0, 0, width=SLIDE_W, height=SLIDE_H)
        veil = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        _no_line(veil)
        veil.fill.solid()
        veil.fill.fore_color.rgb = _rgb(self.t["bg"])
        _alpha(veil, scrim)

    def _accent_bar(self, slide, top, height=Inches(1.6)):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top,
                                     Inches(0.09), height)
        _no_line(bar)
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(self.t["accent"])

    def _footer(self, slide, light=True):
        """Altbet — HAR DOIM band joyda, hech qachon kontent ustida emas."""
        self.n += 1
        tf = _textbox(slide, MARGIN, SLIDE_H - FOOTER_H,
                      SLIDE_W - 2 * MARGIN, Inches(0.3))
        color = self.t["muted"] if light else "8A97A5"
        _para(tf, f"{self.footer_text}  |  {self.n}", 10, color, first=True)

    def _credit(self, slide, left, top, width, text, light=True,
                align=PP_ALIGN.LEFT):
        """Rasm manbasi — rasm QUTISIDAN TASHQARIDA, ustida emas.

        Ilgari model buni rasm ustiga qo'yib yuborardi va yozuv rasmga
        aralashib ketardi. Bu yerda joy oldindan ajratilgan.
        """
        if not text:
            return
        tf = _textbox(slide, left, top, width, Inches(0.24))
        _para(tf, f"Manba: {text}", 9,
              self.t["muted"] if light else "9AA5B1", align=align, first=True)

    def _title(self, slide, text, top=None, size=34):
        top = MARGIN if top is None else top
        tf = _textbox(slide, MARGIN + Inches(0.28), top,
                      SLIDE_W - 2 * MARGIN - Inches(0.28), Inches(0.9))
        _para(tf, text, size, self.t["fg"], bold=True, first=True)
        self._accent_bar(slide, top + Inches(0.06), Inches(0.62))
        return top + Inches(1.0)

    # ── MAKETLAR ────────────────────────────────────────────────────
    def cover(self, title, subtitle=None, image=AUTO, credit=None, footer=None):
        """Muqova. Rasm berilsa — TO'LIQ EKRANLI fon (eng ta'sirli ko'rinish)."""
        s = self._blank()
        image = self._pick(image)
        if image:
            self._bg_photo(s, image, scrim=70)
        else:
            self._bg_solid(s)

        block_top = Inches(2.5)
        self._accent_bar(s, block_top, Inches(2.0))
        tf = _textbox(s, MARGIN + Inches(0.35), block_top,
                      SLIDE_W - 2 * MARGIN - Inches(0.35), Inches(1.5))
        size = 54 if len(str(title)) <= 34 else 40
        _para(tf, title, size, self.t["fg"], bold=True, space_after=10, first=True)
        if subtitle:
            _para(tf, subtitle, 20, self.t["muted"])
        if footer:
            tf2 = _textbox(s, MARGIN + Inches(0.35), Inches(5.6),
                           SLIDE_W - 2 * MARGIN, Inches(0.4))
            _para(tf2, footer, 14, self.t["accent"], bold=True, first=True)
        if credit:
            self._credit(s, MARGIN, SLIDE_H - FOOTER_H - Inches(0.28),
                         Inches(5), credit)
        self._footer(s)
        return s

    def section(self, title, number=None, image=AUTO):
        """Bo'lim ajratuvchi — katta raqam va sarlavha."""
        s = self._blank()
        image = self._pick(image)
        if image:
            self._bg_photo(s, image, scrim=80)
        else:
            self._bg_solid(s)
        top = Inches(2.7)
        if number is not None:
            # ⚠️ Quti balandligi sarlavha boshlanishidan OLDIN tugashi shart
            # (1.65 + 0.9 = 2.55 < 2.7). Ilgari 0.1 dyuymga kesishardi —
            # ko'zga tashlanmasdi, lekin aynan shunday mayda kesishishlar
            # yig'ilib slaydni "dabdala" qiladi.
            tf = _textbox(s, MARGIN, top - Inches(1.05), Inches(3), Inches(0.9))
            _para(tf, f"{int(number):02d}", 64, self.t["accent"],
                  bold=True, first=True)
        self._accent_bar(s, top, Inches(0.9))
        tf = _textbox(s, MARGIN + Inches(0.35), top,
                      SLIDE_W - 2 * MARGIN - Inches(0.35), Inches(1.2))
        _para(tf, title, 40, self.t["fg"], bold=True, first=True)
        self._footer(s)
        return s

    def bullets(self, title, items, subtitle=None, image=AUTO, credit=None):
        """Asosiy maket: chapda matn, o'ngda rasm (rasm bo'lmasa — keng matn).

        `items` — satrlar ro'yxati yoki ("Sarlavha", "Izoh") juftliklari.
        HAR BAND ALOHIDA XAT BOSHI bo'ladi, ular bir-biriga tiqilmaydi.

        ⚠️ Matn ham, rasm ham BIR XIL balandlikdagi, bir xil yumaloq
        burchakli ramkada turadi va bir qatorda boshlanadi — shuning
        uchun slayd muvozanatli ko'rinadi, rasm esa "goh katta, goh
        kichik" bo'lib qolmaydi.
        """
        s = self._blank()
        self._bg_solid(s)
        body_top = self._title(s, title)
        image = self._pick(image)

        # Ikkala ustun uchun BIR XIL ramka balandligi. Manba yozuvi uchun
        # joy oldindan ayirib qo'yiladi — u ramka ICHIGA tushmaydi.
        frame_h = SLIDE_H - body_top - FOOTER_H - Inches(0.34)
        pad = Inches(0.32)

        if image:
            col_w = int((SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2)
            img_l = MARGIN + col_w + Inches(0.4)
            self._photo(s, image, img_l, body_top, col_w, frame_h, mode="fill")
            self._credit(s, img_l, body_top + frame_h + Inches(0.06),
                         col_w, credit)
            text_l, text_w = MARGIN, col_w
        else:
            text_l, text_w = MARGIN, SLIDE_W - 2 * MARGIN

        # Matn ham kartochka ichida — rasm bilan bir xil "og'irlik".
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  text_l, body_top, text_w, frame_h)
        _no_line(card)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(self.t["card"])

        inner_l = text_l + pad
        inner_w = text_w - 2 * pad
        inner_top = body_top + pad
        if subtitle:
            tf = _textbox(s, inner_l, inner_top, inner_w, Inches(0.4))
            _para(tf, subtitle, 15, self.t["muted"], first=True)
            inner_top += Inches(0.5)

        items = list(items or [])
        # Shrift band soniga qarab tanlanadi — 6 ta band ham sig'sin.
        size = 20 if len(items) <= 4 else (17 if len(items) <= 6 else 14)
        if image:
            size -= 2                     # ustun ikki barobar tor
        tf = _textbox(s, inner_l, inner_top, inner_w,
                      body_top + frame_h - inner_top - pad)
        for i, item in enumerate(items):
            if isinstance(item, (tuple, list)) and len(item) >= 2:
                head, note = item[0], item[1]
                _para(tf, f"▸  {head}", size, self.t["accent"], bold=True,
                      space_after=2, first=(i == 0))
                _para(tf, f"     {note}", size - 3, self.t["fg"], space_after=10)
            else:
                _para(tf, f"▸  {item}", size, self.t["fg"], space_after=10,
                      first=(i == 0))
        self._footer(s)
        return s

    def image_slide(self, title, image=AUTO, caption=None, credit=None):
        """Katta rasm slaydi. Diagramma (matplotlib PNG) uchun ham shu.

        Bu yerda rasm QIRQILMAYDI (mode="fit"): asosiy fotosurat yoki
        diagrammani qirqish mazmunni yo'qotadi. Ramka o'lchami esa
        doimiy, shuning uchun slayddan slaydga "sakramaydi".
        """
        s = self._blank()
        self._bg_solid(s)
        body_top = self._title(s, title)
        image = self._pick(image)
        # Izoh va manba uchun joy ANIQ hisoblanadi. Ilgari bu taxminiy
        # qiymat edi va ikkalasi birga bo'lganda manba altbet bandiga
        # 0.02 dyuymga kirib ketardi — ko'zga tashlanmaydigan, lekin
        # "dabdala" ko'rinishni yig'adigan xato.
        reserved = Inches(0.14)
        if caption:
            reserved += Inches(0.32)
        if credit:
            reserved += Inches(0.26)
        frame_h = SLIDE_H - body_top - FOOTER_H - reserved
        rect = self._photo(s, image, MARGIN, body_top,
                           SLIDE_W - 2 * MARGIN, frame_h, mode="fit")
        # Izoh va manba RAMKANING haqiqiy ostiga — u kichraygan bo'lsa
        # ular ham yuqoriga suriladi, oradagi bo'shliq ochilib qolmaydi.
        below = rect[1] + rect[3] + Inches(0.10)
        if caption:
            tf = _textbox(s, MARGIN, below, SLIDE_W - 2 * MARGIN, Inches(0.3))
            _para(tf, caption, 13, self.t["fg"], align=PP_ALIGN.CENTER, first=True)
            below += Inches(0.32)
        self._credit(s, MARGIN, below, SLIDE_W - 2 * MARGIN, credit,
                     align=PP_ALIGN.CENTER)
        self._footer(s)
        return s

    def stats(self, title, pairs, note=None):
        """Katta raqamlar qatori — 2 tadan 4 tagacha."""
        s = self._blank()
        self._bg_solid(s)
        body_top = self._title(s, title)
        pairs = list(pairs or [])[:4]
        if not pairs:
            self._footer(s)
            return s
        gap = Inches(0.35)
        card_w = int((SLIDE_W - 2 * MARGIN - gap * (len(pairs) - 1)) / len(pairs))
        card_h = Inches(2.4)
        top = body_top + Inches(0.5)
        for i, item in enumerate(pairs):
            value, label = (list(item) + [""])[:2]
            left = MARGIN + i * (card_w + gap)
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, top, card_w, card_h)
            _no_line(card)
            card.fill.solid()
            card.fill.fore_color.rgb = _rgb(self.t["card"])
            card.text_frame.text = ""
            tf = _textbox(s, left + Inches(0.25), top + Inches(0.45),
                          card_w - Inches(0.5), card_h - Inches(0.7))
            _para(tf, value, 40, self.t["accent"], bold=True,
                  align=PP_ALIGN.CENTER, space_after=4, first=True)
            _para(tf, label, 14, self.t["muted"], align=PP_ALIGN.CENTER)
        if note:
            tf = _textbox(s, MARGIN, top + card_h + Inches(0.35),
                          SLIDE_W - 2 * MARGIN, Inches(0.4))
            _para(tf, note, 13, self.t["muted"], first=True)
        self._footer(s)
        return s

    def table(self, title, rows, note=None):
        """Jadval. `rows[0]` — sarlavha qatori."""
        s = self._blank()
        self._bg_solid(s)
        body_top = self._title(s, title)
        rows = [list(r) for r in (rows or []) if r]
        if not rows:
            self._footer(s)
            return s
        ncols = max(len(r) for r in rows)
        rows = [r + [""] * (ncols - len(r)) for r in rows]
        height = min(Inches(0.45) * len(rows),
                     SLIDE_H - body_top - FOOTER_H - Inches(0.3))
        shape = s.shapes.add_table(len(rows), ncols, MARGIN, body_top,
                                   SLIDE_W - 2 * MARGIN, height)
        tbl = shape.table
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(
                    self.t["accent"] if ri == 0 else self.t["card"])
                p = cell.text_frame.paragraphs[0]
                p.font.name, p.font.size = FONT, Pt(13 if ri else 14)
                p.font.bold = ri == 0
                p.font.color.rgb = _rgb(self.t["bg"] if ri == 0 else self.t["fg"])
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if note:
            tf = _textbox(s, MARGIN, body_top + height + Inches(0.2),
                          SLIDE_W - 2 * MARGIN, Inches(0.4))
            _para(tf, note, 12, self.t["muted"], first=True)
        self._footer(s)
        return s

    def quote(self, text, author=None, image=AUTO):
        """Kuchli iqtibos — bo'limlar orasida ritm beradi."""
        s = self._blank()
        image = self._pick(image)
        if image:
            self._bg_photo(s, image, scrim=82)
        else:
            self._bg_solid(s)
        tf = _textbox(s, MARGIN + Inches(0.8), Inches(2.4),
                      SLIDE_W - 2 * MARGIN - Inches(1.6), Inches(2.4))
        _para(tf, f"“{text}”", 28, self.t["fg"], bold=True,
              align=PP_ALIGN.CENTER, space_after=14, first=True)
        if author:
            _para(tf, f"— {author}", 16, self.t["accent"], align=PP_ALIGN.CENTER)
        self._footer(s)
        return s

    def closing(self, title="Xulosa", subtitle=None, image=AUTO):
        s = self._blank()
        image = self._pick(image)
        if image:
            self._bg_photo(s, image, scrim=78)
        else:
            self._bg_solid(s)
        self._accent_bar(s, Inches(2.9), Inches(1.4))
        tf = _textbox(s, MARGIN + Inches(0.35), Inches(2.9),
                      SLIDE_W - 2 * MARGIN - Inches(0.35), Inches(1.6))
        _para(tf, title, 44, self.t["fg"], bold=True, space_after=10, first=True)
        if subtitle:
            _para(tf, subtitle, 19, self.t["muted"])
        self._footer(s)
        return s

    # ── SAQLASH ─────────────────────────────────────────────────────
    def _add_transitions(self):
        """Har slaydga yumshoq "fade" o'tish qo'shadi.

        PPTX'da haqiqiy animatsiya (<p:timing>) juda murakkab va mo'rt,
        lekin SLAYD O'TISHI — bitta kichik XML bloki. Taqdimot ko'rsatilganda
        sezilarli darajada jonli ko'rinadi.

        Xato bo'lsa jimgina o'tkazib yuboriladi — o'tish bezak, fayl esa
        undan muhimroq.
        """
        xml = (
            '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org'
            '/markup-compatibility/2006" xmlns:p="http://schemas.openxmlformats'
            '.org/presentationml/2006/main">'
            '<mc:Choice xmlns:p14="http://schemas.microsoft.com/office/'
            'powerpoint/2010/main" Requires="p14">'
            '<p:transition spd="med" p14:dur="600"><p:fade/></p:transition>'
            '</mc:Choice>'
            '<mc:Fallback><p:transition spd="med"><p:fade/></p:transition>'
            '</mc:Fallback></mc:AlternateContent>'
        )
        for slide in self.prs.slides:
            try:
                sld = slide._element
                # Sxema tartibi: cSld -> clrMapOvr -> transition -> timing.
                anchor = sld.find(qn("p:clrMapOvr"))
                node = parse_xml(xml)
                if anchor is not None:
                    anchor.addnext(node)
                else:
                    sld.append(node)
            except Exception:
                pass

    def save(self, path):
        """Faylni saqlaydi. Papka yo'q bo'lsa yaratadi."""
        self._add_transitions()
        folder = os.path.dirname(path)
        if folder:
            os.makedirs(folder, exist_ok=True)
        self.prs.save(path)
        return path


def demo(path="output/demo.pptx"):
    """Modulning o'zini tekshirish: `python deck.py`."""
    d = Deck("Demo taqdimot", theme="navy", footer="Demo")
    d.cover("KATTA SARLAVHA", "Kichik izoh matni", footer="Fan nomi")
    d.section("Birinchi bo'lim", 1)
    d.bullets("Asosiy fikrlar",
              [("Birinchi", "Izoh matni"), ("Ikkinchi", "Yana izoh"),
               ("Uchinchi", "Uchinchi izoh")])
    d.stats("Raqamlar", [("38 mln", "Talafot"), ("4", "Yil"), ("30+", "Davlat")])
    d.table("Solishtirish", [["Nom", "Qiymat"], ["A", "1"], ["B", "2"]])
    d.quote("Qisqa va kuchli iqtibos.", "Muallif")
    d.closing("Xulosa", "Yakuniy fikr")
    return d.save(path)


if __name__ == "__main__":
    out = demo()
    assert os.path.exists(out) and os.path.getsize(out) > 20000, "PPTX yaratilmadi"
    # Qayta ochilishi — XML buzilmaganini tekshiradi (o'tishlar qo'shilgach).
    reopened = Presentation(out)
    assert len(reopened.slides.__iter__.__self__._sldIdLst) == 7, "slayd soni xato"
    print(f"[1] demo() {len(reopened.slides.__iter__.__self__._sldIdLst)} slayd OK")
    print(f"[2] fayl qayta ochildi ({os.path.getsize(out) // 1024} KB) OK")
    print("✅ deck.py ishlayapti")
