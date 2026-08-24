"""deck.py maketlarining GEOMETRIYASI uchun tekshiruvlar.

Ishga tushirish:
    PYTHONIOENCODING=utf-8 python tests/test_deck_layout.py

NEGA KERAK: jonli sinovda slaydlar "dabdala" chiqqan edi — rasm matn
ustiga, manba yozuvi rasm ustiga, altbet raqami esa rasm ustiga tushgan.
Bu yerdagi tekshiruvlar aynan shu uchtasini qaytib kelishiga yo'l qo'ymaydi.
Tarmoqqa chiqmaydi — rasm mahalliy yaratiladi.
"""
import io
import os
import sys

sys.path.insert(0, os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "services", "sandbox_helpers"))

from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

import deck

failures = []


def check(n, label, cond):
    print(f"[{n}] {label} {'OK' if cond else 'XATO'}")
    if not cond:
        failures.append(label)


TMP = os.environ.get("TEMP", ".")
PORTRAIT = os.path.join(TMP, "_deck_portrait.jpg")
WIDE = os.path.join(TMP, "_deck_wide.jpg")
Image.new("RGB", (600, 1000), (120, 60, 40)).save(PORTRAIT)
Image.new("RGB", (2000, 800), (40, 90, 120)).save(WIDE)


def boxes(slide):
    """Slayddagi KONTENT shakllari: fon to'ldiruvchilari hisobga olinmaydi."""
    out = []
    for sh in slide.shapes:
        w, h = sh.width or 0, sh.height or 0
        full_bleed = w >= deck.SLIDE_W * 0.99 and h >= deck.SLIDE_H * 0.99
        if full_bleed:
            continue                      # fon to'rtburchagi / fon rasmi
        out.append((sh, sh.left or 0, sh.top or 0, w, h))
    return out


def overlap(a, b):
    _, al, at, aw, ah = a
    _, bl, bt, bw, bh = b
    return not (al + aw <= bl or bl + bw <= al or at + ah <= bt or bt + bh <= at)


# ─────────────────────────────────────────────────────────────
# 1-2. RASM NISBATI
# ─────────────────────────────────────────────────────────────
l, t, w, h = deck._fit_box(PORTRAIT, 0, 0, deck.Inches(4), deck.Inches(4))
check(1, "portret rasm nisbati saqlanadi (cho'zilmaydi)",
      abs((w / h) - (600 / 1000)) < 0.02 and w <= deck.Inches(4) + 1
      and h <= deck.Inches(4) + 1)

l, t, w, h = deck._fit_box(WIDE, 0, 0, deck.Inches(4), deck.Inches(4))
check(2, "keng rasm ham qutiga sig'adi va markazlanadi",
      abs((w / h) - 2.5) < 0.05 and t > 0)

cropped = deck._cover_crop(PORTRAIT, deck.SLIDE_W, deck.SLIDE_H)
with Image.open(cropped) as im:
    cw, ch = im.size
check(3, "fon rasmi 16:9 ga qirqiladi (yon tomonda bo'sh joy qolmaydi)",
      abs((cw / ch) - (deck.SLIDE_W / deck.SLIDE_H)) < 0.02)


# ─────────────────────────────────────────────────────────────
# 4-7. MAKETLAR: CHEGARA, USTMA-USTLIK, ALTBET
# ─────────────────────────────────────────────────────────────
d = deck.Deck("Test taqdimot", theme="navy", footer="Test")
d.cover("BIRINCHI JAHON URUSHI", "1914–1918", image=WIDE,
        credit="iwm.org.uk", footer="Tarix fanidan")
d.section("Urush sabablari", 1)
d.bullets("Asosiy sabablar",
          [("Imperializm", "Mustamlakalar uchun kurash"),
           ("Militarizm", "Qurollanish poygasi"),
           ("Ittifoqlar", "Ikki blokka bo'linish")],
          image=PORTRAIT, credit="wikimedia.org")
d.image_slide("Xandaq urushi", PORTRAIT, "G'arbiy front, 1916",
              credit="iwm.org.uk")
d.stats("Raqamlarda", [("38 mln", "Talafot"), ("4", "Yil"), ("30+", "Davlat")])
d.table("Solishtirish", [["Davlat", "Talafot"], ["Rossiya", "1.8 mln"]])
d.closing("Xulosa", "Urush jahon tartibini o'zgartirdi")
out = d.save(os.path.join(TMP, "_deck_test.pptx"))

from pptx import Presentation
prs = Presentation(out)
slides = list(prs.slides)

inside = True
for s in slides:
    for sh, sl, st, sw, sh_ in boxes(s):
        if sl < -1000 or st < -1000 or sl + sw > deck.SLIDE_W + 1000 \
                or st + sh_ > deck.SLIDE_H + 1000:
            inside = False
check(4, "hech bir element slayd chegarasidan chiqmaydi", inside)

clashes = []
for i, s in enumerate(slides):
    bs = boxes(s)
    for a in range(len(bs)):
        for b in range(a + 1, len(bs)):
            # Kartochka (stats) ichidagi matn ataylab kartochka ustida —
            # bu maket, ustma-ustlik emas.
            sa, sb = bs[a][0], bs[b][0]
            if sa.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or \
               sb.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if overlap(bs[a], bs[b]):
                clashes.append((i + 1, sa.shape_type, sb.shape_type))
check(5, "matn va rasm bir-birining USTIGA chiqmaydi", not clashes)
if clashes:
    print("     to'qnashuvlar:", clashes[:5])

# Altbet bandiga faqat altbet matni tushishi kerak.
footer_top = deck.SLIDE_H - deck.FOOTER_H
intruders = []
for i, s in enumerate(slides):
    for sh, sl, st, sw, sh_ in boxes(s):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and st + sh_ > footer_top:
            intruders.append(i + 1)
check(6, "rasm altbet bandiga kirmaydi (raqam rasm ustida qolmaydi)",
      not intruders)
if intruders:
    print("     buzgan slaydlar:", intruders)

# Manba yozuvi rasmning ostida bo'lishi kerak, ustida emas.
bullets_slide = slides[2]
pics = [b for b in boxes(bullets_slide)
        if b[0].shape_type == MSO_SHAPE_TYPE.PICTURE]
credits = [b for b in boxes(bullets_slide)
           if b[0].has_text_frame and "Manba:" in b[0].text_frame.text]
check(7, "manba yozuvi rasm QUTISIDAN TASHQARIDA (ostida)",
      len(pics) == 1 and len(credits) == 1
      and credits[0][2] >= pics[0][2] + pics[0][4] - 1000)


# ─────────────────────────────────────────────────────────────
# 8-10. MAZMUN VA MUSTAHKAMLIK
# ─────────────────────────────────────────────────────────────
check(8, "har slaydda altbet raqami bor",
      all(any(b[0].has_text_frame and "|" in b[0].text_frame.text
              for b in boxes(s)) for s in slides))

# Rasm yo'q bo'lsa ham yiqilmasin (rasm topilmagan holat).
d2 = deck.Deck("Rasmsiz")
d2.cover("Sarlavha", "Izoh", image="yoq_bunday_fayl.jpg")
d2.bullets("Matn", ["Bir", "Ikki"], image="yoq_bunday_fayl.jpg", credit="x.com")
d2.image_slide("Rasm", "yoq_bunday_fayl.jpg", "izoh")
out2 = d2.save(os.path.join(TMP, "_deck_test2.pptx"))
check(9, "mavjud bo'lmagan rasm bilan ham yiqilmaydi",
      os.path.exists(out2) and len(list(Presentation(out2).slides)) == 3)

check(10, "o'tish (transition) qo'shilgach fayl qayta ochiladi",
      len(list(Presentation(out).slides)) == 7 and os.path.getsize(out) > 20000)


# ─────────────────────────────────────────────────────────────
# 11-14. RASM HOVUZI VA RAMKA MUVOZANATI
# ─────────────────────────────────────────────────────────────
# Model odatda rasmni FAQAT muqovaga yozib, qolgan slaydlarni bo'sh
# qoldiradi. Hovuz shuning uchun bor: maketlar rasmni O'ZI oladi.
d3 = deck.Deck("Hovuz", images=[WIDE, PORTRAIT])
d3.cover("Muqova", "izoh")
d3.section("Bo'lim", 1)
d3.bullets("Matn", [("A", "izoh"), ("B", "izoh")])
d3.closing("Xulosa")
out3 = d3.save(os.path.join(TMP, "_deck_pool.pptx"))
pool_slides = list(Presentation(out3).slides)


def pic_count(slide):
    return sum(1 for sh in slide.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)


check(11, "rasm faqat muqovada emas — hamma mos slaydga tarqaladi",
      all(pic_count(s) >= 1 for s in pool_slides))

d4 = deck.Deck("Rasmsiz", images=[WIDE])
d4.bullets("Faqat matn", ["Bir", "Ikki"], image=None)
out4 = d4.save(os.path.join(TMP, "_deck_none.pptx"))
check(12, "image=None berilsa rasm ATAYLAB qo'yilmaydi",
      pic_count(list(Presentation(out4).slides)[0]) == 0)

# Matn kartochkasi va rasm ramkasi bir qatorda, bir xil balandlikda.
bs = boxes(slides[2])
pic = [b for b in bs if b[0].shape_type == MSO_SHAPE_TYPE.PICTURE][0]
cards = [b for b in bs if b[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
         and b[3] > deck.Inches(3)]
check(13, "matn kartochkasi va rasm ramkasi bir xil o'lchamda va bir qatorda",
      cards and abs(cards[0][2] - pic[2]) < 20000
      and abs(cards[0][4] - pic[4]) < 20000
      and abs(cards[0][3] - pic[3]) < 20000)

# Rasm ramkani TO'LIQ to'ldiradi (fill) — yon tomonda bo'sh joy qolmaydi.
check(14, "yon ustundagi rasm ramkani to'liq to'ldiradi",
      abs(pic[3] - cards[0][3]) < 20000 and abs(pic[4] - cards[0][4]) < 20000)


# ─────────────────────────────────────────────────────────────
# 15-16. RAMKA RASMGA YOPISHADI (bo'sh joy qolmaydi)
# ─────────────────────────────────────────────────────────────
# Xarita yoki tik rasm keng ramka ichida kichkina bo'lib qolib, yon
# tomonlarda katta bo'sh maydon qolardi.
img_slide = slides[3]                     # image_slide (PORTRAIT rasm)
ib = boxes(img_slide)
ipic = [b for b in ib if b[0].shape_type == MSO_SHAPE_TYPE.PICTURE][0]
icard = [b for b in ib if b[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
         and b[3] > deck.Inches(1)][0]
slack_w = icard[3] - ipic[3]
slack_h = icard[4] - ipic[4]
check(15, "image_slide: ramka rasmni QUCHOQLAYDI (ortiqcha bo'sh joy yo'q)",
      0 <= slack_w <= deck.Inches(0.4) and 0 <= slack_h <= deck.Inches(0.4))
if not (0 <= slack_w <= deck.Inches(0.4)):
    print(f"     ortiqcha kenglik: {slack_w / 914400:.2f} dyuym")

# Izoh ramkaning ostida va unga yaqin turishi kerak.
caps = [b for b in ib if b[0].has_text_frame
        and "front" in b[0].text_frame.text]
check(16, "izoh ramkaning ANIQ ostida (oradagi bo'shliq ochilmaydi)",
      caps and 0 < caps[0][2] - (icard[2] + icard[4]) < deck.Inches(0.35))


print("─" * 55)
if failures:
    print(f"❌ {len(failures)} ta tekshiruv yiqildi:")
    for f in failures:
        print(f"   • {f}")
    sys.exit(1)
print("✅ deck.py maket tekshiruvlari — hammasi o'tdi")
